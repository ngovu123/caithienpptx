"""
A set of functions to create a PowerPoint slide deck.
"""
import logging
import os
import pathlib
import random
import re
import sys
import tempfile
from typing import List, Tuple, Optional

import json5
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

sys.path.append('..')
sys.path.append('../..')

import helpers.image_search as ims
from global_config import GlobalConfig


load_dotenv()


# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'

ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1

IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")

ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]


logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)


def remove_slide_number_from_heading(header: str) -> str:
    """
    Remove the slide number from a given slide header.

    :param header: The header of a slide.
    :return: The header without slide number.
    """

    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:]

    return header


def generate_powerpoint_presentation(
        parsed_data: dict,
        slides_template: str,
        output_file_path: pathlib.Path
) -> List:
    """
    Create and save a PowerPoint presentation file containing the content in JSON format.

    :param parsed_data: The presentation content as parsed JSON data.
    :param slides_template: The PPTX template to use.
    :param output_file_path: The path of the PPTX file to save as.
    :return: A list of presentation title and slides headers.
    """

    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(presentation)

    # The title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = parsed_data['title']
    logger.info(
        'PPT title: %s | #slides: %d | template: %s',
        title.text, len(parsed_data['slides']),
        GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file']
    )
    subtitle.text = 'bởi tôi và SparkHub AI :)'
    all_headers = [title.text, ]

    # Add content in a loop
    for a_slide in parsed_data['slides']:
        is_processing_done = _handle_icons_ideas(
            presentation=presentation,
            slide_json=a_slide,
            slide_width_inch=slide_width_inch,
            slide_height_inch=slide_height_inch
        )

        if not is_processing_done:
            is_processing_done = _handle_double_col_layout(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

        if not is_processing_done:
            is_processing_done = _handle_step_by_step_process(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

        if not is_processing_done:
            _handle_default_display(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

    # The thank-you slide
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    title = slide.shapes.title
    title.text = 'Thank you!'

    presentation.save(output_file_path)

    return all_headers


def get_flat_list_of_contents(items: list, level: int) -> List[Tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and
    its level.

    :param items: A bullet point (string or list).
    :param level: The current level of hierarchy.
    :return: A list of (bullet item text, hierarchical level) tuples.
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list


def get_slide_placeholders(
        slide: pptx.slide.Slide,
        layout_number: int,
        is_debug: bool = False
) -> List[Tuple[int, str]]:
    """
    Return the index and name (lower case) of all placeholders present in a slide, except
    the title placeholder.

    A placeholder in a slide is a place to add content. Each placeholder has a name and an index.
    This index is NOT a list index, rather a set of keys used to look up a dict. So, `idx` is
    non-contiguous. Also, the title placeholder of a slide always has index 0. User-added
    placeholder get indices assigned starting from 10.

    With user-edited or added placeholders, their index may be difficult to track. This function
    returns the placeholders name as well, which could be useful to distinguish between the
    different placeholder.

    :param slide: The slide.
    :param layout_number: The layout number used by the slide.
    :param is_debug: Whether to print debugging statements.
    :return: A list containing placeholders (idx, name) tuples, except the title placeholder.
    """

    if is_debug:
        print(
            f'Slide layout #{layout_number}:'
            f' # of placeholders: {len(slide.shapes.placeholders)} (including the title)'
        )

    placeholders = [
        (shape.placeholder_format.idx, shape.name.lower()) for shape in slide.shapes.placeholders
    ]
    placeholders.pop(0)  # Remove the title placeholder

    if is_debug:
        print(placeholders)

    return placeholders


def _handle_default_display(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Display a list of text in a slide.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    """

    status = False

    if 'img_keywords' in slide_json:
        if random.random() < IMAGE_DISPLAY_PROBABILITY:
            if random.random() < FOREGROUND_IMAGE_PROBABILITY:
                status = _handle_display_image__in_foreground(
                    presentation,
                    slide_json,
                    slide_width_inch,
                    slide_height_inch
                )
            else:
                status = _handle_display_image__in_background(
                    presentation,
                    slide_json,
                    slide_width_inch,
                    slide_height_inch
                )

    if status:
        return

    # Image display failed, so display only text
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title

    try:
        body_shape = shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        body_shape = shapes.placeholders[placeholders[0][0]]

    title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    text_frame = body_shape.text_frame

    # The bullet_points may contain a nested hierarchy of JSON arrays
    # In some scenarios, it may contain objects (dictionaries) because the LLM generated so
    #  ^ The second scenario is not covered

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    _handle_key_message(
        the_slide=slide,
        slide_json=slide_json,
        slide_height_inch=slide_height_inch,
        slide_width_inch=slide_width_inch
    )


def _handle_display_image__in_foreground(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Create a slide with text and image using a picture placeholder layout. If not image keyword is
    available, it will add only text to the slide.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if the side has been processed.
    """

    img_keywords = slide_json['img_keywords'].strip()
    slide = presentation.slide_layouts[8]  # Picture with Caption
    slide = presentation.slides.add_slide(slide)
    placeholders = None

    title_placeholder = slide.shapes.title
    title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])

    try:
        pic_col: PicturePlaceholder = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=8)
        pic_col = None
        for idx, name in placeholders:
            if 'picture' in name:
                pic_col: PicturePlaceholder = slide.shapes.placeholders[idx]

    try:
        text_col: SlidePlaceholder = slide.shapes.placeholders[2]
    except KeyError:
        text_col = None
        if not placeholders:
            placeholders = get_slide_placeholders(slide, layout_number=8)

        for idx, name in placeholders:
            if 'content' in name:
                text_col: SlidePlaceholder = slide.shapes.placeholders[idx]

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            text_col.text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = text_col.text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    if not img_keywords:
        # No keywords, so no image search and addition
        return True

    try:
        photo_url, page_url = ims.get_photo_url_from_api_response(
            ims.search_pexels(query=img_keywords, size='medium')
        )

        if photo_url:
            pic_col.insert_picture(
                ims.get_image_from_url(photo_url)
            )

            _add_text_at_bottom(
                slide=slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch,
                text='Photo provided by Pexels',
                hyperlink=page_url
            )
    except Exception as ex:
        logger.error(
            '*** Error occurred while running adding image to slide: %s',
            str(ex)
        )

    return True


def _handle_display_image__in_background(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with text and an image in the background. It works just like
    `_handle_default_display()` but with a background image added. If not image keyword is
    available, it will add only text to the slide.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if the slide has been processed.
    """

    img_keywords = slide_json['img_keywords'].strip()

    # Add a photo in the background, text in the foreground
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title

    try:
        body_shape = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        # Layout 1 usually has two placeholders, including the title
        body_shape = slide.shapes.placeholders[placeholders[0][0]]

    title_shape.text = remove_slide_number_from_heading(slide_json['heading'])

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            body_shape.text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    if not img_keywords:
        # No keywords, so no image search and addition
        return True

    try:
        photo_url, page_url = ims.get_photo_url_from_api_response(
            ims.search_pexels(query=img_keywords, size='large')
        )

        if photo_url:
            picture = slide.shapes.add_picture(
                image_file=ims.get_image_from_url(photo_url),
                left=0,
                top=0,
                width=pptx.util.Inches(slide_width_inch),
            )

            _add_text_at_bottom(
                slide=slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch,
                text='Photo provided by Pexels',
                hyperlink=page_url
            )

            # Move picture to background
            # https://github.com/scanny/python-pptx/issues/49#issuecomment-137172836
            slide.shapes._spTree.remove(picture._element)
            slide.shapes._spTree.insert(2, picture._element)
    except Exception as ex:
        logger.error(
            '*** Error occurred while running adding image to the slide background: %s',
            str(ex)
        )

    return True


def _handle_icons_ideas(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Add a slide with some icons and text.
    If no suitable icons are found, the step numbers are shown.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if the slide has been processed.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        items = slide_json['bullet_points']

        # Ensure that it is a single list of strings without any sub-list
        for step in items:
            if not isinstance(step, str) or not step.startswith(ICON_BEGINNING_MARKER):
                return False

        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = remove_slide_number_from_heading(slide_json['heading'])

        n_items = len(items)
        text_box_size = INCHES_2

        # Calculate the total width of all pictures and the spacing
        total_width = n_items * ICON_SIZE
        # slide_width = presentation.slide_width
        spacing = (pptx.util.Inches(slide_width_inch) - total_width) / (n_items + 1)

        for idx, item in enumerate(items):
            # Extract the icon name and text
            match = ICONS_REGEX.search(item)

            if not match:
                # print('No icon/text pattern match found...skipping to the next item')
                continue

            icon_name = match.group(1)
            accompanying_text = match.group(2)
            icon_path = f'{GlobalConfig.ICONS_DIR}/{icon_name}.png'

            left = spacing + idx * (ICON_SIZE + spacing)
            top = INCHES_3

            # Calculate the center position for alignment
            center = left + ICON_SIZE / 2

            # Add a rectangle shape with a fill color (background)
            # The size of the shape is slightly bigger than the icon, so align the icon position
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                center - INCHES_0_5,
                top - (ICON_BG_SIZE - ICON_SIZE) / 2,
                INCHES_1, INCHES_1
            )
            shape.fill.solid()
            shape.shadow.inherit = False

            # Set the icon's background shape color
            color = random.choice(ICON_COLORS)
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = color

            # Add the icon image on top of the colored shape
            try:
                slide.shapes.add_picture(icon_path, left, top, height=ICON_SIZE)
            except FileNotFoundError:
                logger.error(
                    'Icon %s not found...using generic step number as icon...',
                    icon_name
                )
                step_icon_path = f'{GlobalConfig.ICONS_DIR}/{idx + 1}-circle.png'
                if os.path.exists(step_icon_path):
                    slide.shapes.add_picture(step_icon_path, left, top, height=ICON_SIZE)

            # Add a text box below the shape
            text_top = top + ICON_SIZE + INCHES_0_2
            text_left = center - text_box_size / 2  # Center the text box horizontally
            # text_box = slide.shapes.add_textbox(text_left, text_top, text_box_size, text_box_size)
            text_box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                text_left, text_top,
                text_box_size, text_box_size
            )
            text_frame = text_box.text_frame
            text_frame.text = accompanying_text
            text_frame.word_wrap = True
            text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER

            # Center the text vertically
            text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            text_box.fill.background()  # No fill
            text_box.line.fill.background()  # No line
            text_box.shadow.inherit = False

            # Set the font color based on the theme
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_2

            _add_text_at_bottom(
                slide=slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch,
                text='More icons available in the SlideDeck AI repository',
                hyperlink='https://github.com/barun-saha/slide-deck-ai/tree/main/icons/png128'
            )

        return True

    return False


def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower left side of a slide.

    :param slide: The slide.
    :param slide_width_inch: The width of the slide.
    :param slide_height_inch: The height of the slide.
    :param target_height: the target height of the box in inches (optional).
    :param text: The text to be added
    :param hyperlink: The hyperlink to be added to the text (optional).
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )

    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = pptx.util.Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink


def _handle_double_col_layout(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with a double column layout for comparison.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if double col layout has been added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        double_col_content = slide_json['bullet_points']

        if double_col_content and (
                len(double_col_content) == 2
        ) and isinstance(double_col_content[0], dict) and isinstance(double_col_content[1], dict):
            slide = presentation.slide_layouts[4]
            slide = presentation.slides.add_slide(slide)
            placeholders = None

            shapes = slide.shapes
            title_placeholder = shapes.title
            title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])

            try:
                left_heading, right_heading = shapes.placeholders[1], shapes.placeholders[3]
            except KeyError:
                # For manually edited/added master slides, the placeholder idx numbers in the dict
                # will be different (>= 10)
                left_heading, right_heading = None, None
                placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'text placeholder' in name:
                        if not left_heading:
                            left_heading = shapes.placeholders[idx]
                        elif not right_heading:
                            right_heading = shapes.placeholders[idx]

            try:
                left_col, right_col = shapes.placeholders[2], shapes.placeholders[4]
            except KeyError:
                left_col, right_col = None, None
                if not placeholders:
                    placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'content placeholder' in name:
                        if not left_col:
                            left_col = shapes.placeholders[idx]
                        elif not right_col:
                            right_col = shapes.placeholders[idx]

            left_col_frame, right_col_frame = left_col.text_frame, right_col.text_frame

            if 'heading' in double_col_content[0] and left_heading:
                left_heading.text = double_col_content[0]['heading']
            if 'bullet_points' in double_col_content[0]:
                flat_items_list = get_flat_list_of_contents(
                    double_col_content[0]['bullet_points'], level=0
                )

                if not left_heading:
                    left_col_frame.text = double_col_content[0]['heading']

                for idx, an_item in enumerate(flat_items_list):
                    if left_heading and idx == 0:
                        left_col_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                    else:
                        paragraph = left_col_frame.add_paragraph()
                        paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                        paragraph.level = an_item[1]

            if 'heading' in double_col_content[1] and right_heading:
                right_heading.text = double_col_content[1]['heading']
            if 'bullet_points' in double_col_content[1]:
                flat_items_list = get_flat_list_of_contents(
                    double_col_content[1]['bullet_points'], level=0
                )

                if not right_heading:
                    right_col_frame.text = double_col_content[1]['heading']

                for idx, an_item in enumerate(flat_items_list):
                    if right_col_frame and idx == 0:
                        right_col_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                    else:
                        paragraph = right_col_frame.add_paragraph()
                        paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                        paragraph.level = an_item[1]

            _handle_key_message(
                the_slide=slide,
                slide_json=slide_json,
                slide_height_inch=slide_height_inch,
                slide_width_inch=slide_width_inch
            )

            return True

    return False


def _handle_step_by_step_process(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add shapes to display a step-by-step process in the slide, if available.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return True if this slide has a step-by-step process depiction added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        steps = slide_json['bullet_points']

        no_marker_count = 0.0
        n_steps = len(steps)

        # Ensure that it is a single list of strings without any sub-list
        for step in steps:
            if not isinstance(step, str):
                return False

            # In some cases, one or two steps may not begin with >>, e.g.:
            # {
            #     "heading": "Step-by-Step Process: Creating a Legacy",
            #     "bullet_points": [
            #         "Identify your unique talents and passions",
            #         ">> Develop your skills and knowledge",
            #         ">> Create meaningful work",
            #         ">> Share your work with the world",
            #         ">> Continuously learn and adapt"
            #     ],
            #     "key_message": ""
            # },
            #
            # Use a threshold, e.g., at most 20%
            if not step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                no_marker_count += 1

        slide_header = slide_json['heading'].lower()
        if (no_marker_count / n_steps > 0.25) and not (
                ('step-by-step' in slide_header) or ('step by step' in slide_header)
        ):
            return False

        if n_steps < 3 or n_steps > 6:
            # Two steps -- probably not a process
            # More than 5--6 steps -- would likely cause a visual clutter
            return False

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = remove_slide_number_from_heading(slide_json['heading'])

        if 3 <= n_steps <= 4:
            # Horizontal display
            height = INCHES_1_5
            width = pptx.util.Inches(slide_width_inch / n_steps - 0.01)
            top = pptx.util.Inches(slide_height_inch / 2)
            left = pptx.util.Inches((slide_width_inch - width.inches * n_steps) / 2 + 0.05)

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
                shape.text = step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                left += width - INCHES_0_4
        elif 4 < n_steps <= 6:
            # Vertical display
            height = pptx.util.Inches(0.65)
            top = pptx.util.Inches(slide_height_inch / 4)
            left = INCHES_1  # slide_width_inch - width.inches)

            # Find the close to median width, based on the length of each text, to be set
            # for the shapes
            width = pptx.util.Inches(slide_width_inch * 2 / 3)
            lengths = [len(step) for step in steps]
            font_size_20pt = pptx.util.Pt(20)
            widths = sorted(
                [
                    min(
                        pptx.util.Inches(font_size_20pt.inches * a_len),
                        width
                    ) for a_len in lengths
                ]
            )
            width = widths[len(widths) // 2]

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height)
                shape.text = step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                top += height + INCHES_0_3
                left += INCHES_0_5

    return True


def _handle_key_message(
        the_slide: pptx.slide.Slide,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Add a shape to display the key message in the slide, if available.

    :param the_slide: The slide to be processed.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    """

    if 'key_message' in slide_json and slide_json['key_message']:
        height = pptx.util.Inches(1.6)
        width = pptx.util.Inches(slide_width_inch / 2.3)
        top = pptx.util.Inches(slide_height_inch - height.inches - 0.1)
        left = pptx.util.Inches((slide_width_inch - width.inches) / 2)
        shape = the_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        shape.text = slide_json['key_message']


def _get_slide_width_height_inches(presentation: pptx.Presentation) -> Tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    :param presentation: The presentation object.
    :return: The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
    # logger.debug('Slide width: %f, height: %f', slide_width_inch, slide_height_inch)

    return slide_width_inch, slide_height_inch


if __name__ == '__main__':
    _JSON_DATA = '''
    {
  "title": "AI Applications: Transforming Industries",
  "slides": [
    {
      "heading": "Introduction to AI Applications",
      "bullet_points": [
        "Artificial Intelligence (AI) is transforming various industries",
        "AI applications range from simple decision-making tools to complex systems",
        "AI can be categorized into types: Rule-based, Instance-based, and Model-based"
      ],
      "key_message": "AI is a broad field with diverse applications and categories",
      "img_keywords": "AI, transformation, industries, decision-making, categories"
    },
    {
      "heading": "AI in Everyday Life",
      "bullet_points": [
        "Virtual assistants like Siri, Alexa, and Google Assistant",
        "Recommender systems in Netflix, Amazon, and Spotify",
        "Fraud detection in banking and credit card transactions"
      ],
      "key_message": "AI is integrated into our daily lives through various services",
      "img_keywords": "virtual assistants, recommender systems, fraud detection"
    },
    {
      "heading": "AI in Healthcare",
      "bullet_points": [
        "Disease diagnosis and prediction using machine learning algorithms",
        "Personalized medicine and drug discovery",
        "AI-powered robotic surgeries and remote patient monitoring"
      ],
      "key_message": "AI is revolutionizing healthcare with improved diagnostics and patient care",
      "img_keywords": "healthcare, disease diagnosis, personalized medicine, robotic surgeries"
    },
    {
      "heading": "AI in Key Industries",
      "bullet_points": [
        {
          "heading": "Retail",
          "bullet_points": [
            "Inventory management and demand forecasting",
            "Customer segmentation and targeted marketing",
            "AI-driven chatbots for customer service"
          ]
        },
        {
          "heading": "Finance",
          "bullet_points": [
            "Credit scoring and risk assessment",
            "Algorithmic trading and portfolio management",
            "AI for detecting money laundering and cyber fraud"
          ]
        }
      ],
      "key_message": "AI is transforming retail and finance with improved operations and decision-making",
      "img_keywords": "retail, finance, inventory management, credit scoring, algorithmic trading"
    },
    {
      "heading": "AI in Education",
      "bullet_points": [
        "Personalized learning paths and adaptive testing",
        "Intelligent tutoring systems for skill development",
        "AI for predicting student performance and dropout rates"
      ],
      "key_message": "AI is personalizing education and improving student outcomes",
    },
    {
      "heading": "Step-by-Step: AI Development Process",
      "bullet_points": [
        ">> Define the problem and objectives",
        ">> Collect and preprocess data",
        ">> Select and train the AI model",
        ">> Evaluate and optimize the model",
        ">> Deploy and monitor the AI system"
      ],
      "key_message": "Developing AI involves a structured process from problem definition to deployment",
      "img_keywords": ""
    },
    {
      "heading": "AI Icons: Key Aspects",
      "bullet_points": [
        "[[brain]] Human-like intelligence and decision-making",
        "[[robot]] Automation and physical tasks",
        "[[cloud]] Data processing and cloud computing",
        "[[lightbulb]] Insights and predictions",
        "[[globe2]] Global connectivity and impact"
      ],
      "key_message": "AI encompasses various aspects, from human-like intelligence to global impact",
      "img_keywords": "AI aspects, intelligence, automation, data processing, global impact"
    },
    {
      "heading": "Conclusion: Embracing AI's Potential",
      "bullet_points": [
        "AI is transforming industries and improving lives",
        "Ethical considerations are crucial for responsible AI development",
        "Invest in AI education and workforce development",
        "Call to action: Explore AI applications and contribute to shaping its future"
      ],
      "key_message": "AI offers immense potential, and we must embrace it responsibly",
      "img_keywords": "AI transformation, ethical considerations, AI education, future of AI"
    }
  ]
}'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(_JSON_DATA),
        output_file_path=path,
        slides_template='Basic'
    )
    print(f'File path: {path}')

    temp.close()